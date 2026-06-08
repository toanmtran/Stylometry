"""
Build the Neural Network slide deck (open-set, 6-class authorship attribution).

Style mirrors the existing K-means / SVM sections of the group deck:
  * warm off-white background, minimalist look
  * short black dash accent above each title
  * Playfair Display serif section titles (running header on the left)
  * Poppins body text, gold accent for key terms / metric chips
  * figure on the right with a bold "Figure N." caption
  * page number bottom-right

Figures come from src/neural_network/outputs/ (regenerate them first with
make_report_figures.py, gradient_diagnostics.py, plot_feature_importance.py).

Output -> docs/neural_network_slides.pptx
"""

from pathlib import Path                          # filesystem paths that work cross-OS

from PIL import Image                             # only used to read each PNG's pixel size (for aspect-ratio fitting)
from pptx import Presentation                     # top-level PowerPoint document object
from pptx.dml.color import RGBColor               # lets us set exact RGB fill/line/font colours
from pptx.enum.shapes import MSO_SHAPE            # shape catalogue (rectangle, rounded rect, arrow, ...)
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN   # vertical text anchoring + horizontal paragraph alignment
from pptx.oxml.ns import qn                       # builds namespaced XML tag names (for low-level font tweaks)
from pptx.util import Inches, Pt                  # unit helpers: Inches/Pt convert to PowerPoint's internal EMUs

NN_DIR = Path(__file__).resolve().parent.parent       # src/neural_network/
FIG = NN_DIR / "outputs"                              # folder holding the pre-generated figures
PROJECT_ROOT = NN_DIR.parents[1]                      # Stylometry/
OUT = PROJECT_ROOT / "docs" / "neural_network_slides.pptx"  # where the finished deck is written

# ---- palette / fonts -------------------------------------------------------
BG = RGBColor(0xEF, 0xEE, 0xEA)      # warm off-white slide background (matches the deck)
INK = RGBColor(0x1A, 0x1A, 0x1A)     # near-black for titles and body text
GOLD = RGBColor(0x94, 0x65, 0x2B)    # brown/gold accent used for bullets, key terms, metric values
GRAY = RGBColor(0x6B, 0x6B, 0x6B)    # muted grey for captions and page numbers
CHIPBG = RGBColor(0xFF, 0xFF, 0xFF)  # white fill for metric "chip" cards and table cells
CHIPLINE = RGBColor(0xDD, 0xDB, 0xD4)  # light grey border for chips / table cells
GREEN = RGBColor(0x3F, 0x8E, 0x4E)   # green accent (softmax box, perfect none-class F1)

SERIF = "Playfair Display"           # deck heading font; used for titles + chip values
SANS = "Poppins"                     # deck body font; used for bullets, captions, labels

SW, SH = 13.333, 7.5                 # slide width/height in inches -> 16:9 widescreen
START_PAGE = 1                       # first page number printed on these slides

prs = Presentation()                 # create a fresh, empty presentation
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]         # layout index 6 is the built-in "Blank" layout (no placeholders)


# ---- low-level helpers -----------------------------------------------------
def add_slide():
    s = prs.slides.add_slide(BLANK)                                              # append a new blank slide
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)  # full-bleed background
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element                                                             # send background rectangle to back
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _set_font(run, font=SANS, size=18, color=INK, bold=False, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()                  # pin typeface in every font slot (avoids substitution)
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0, line=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    return p


def run(p, text, **kw):
    r = p.add_run()
    r.text = text
    _set_font(r, **kw)
    return r


def dash(slide, x=0.66, y=0.6):
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(0.42), Inches(0.055))
    d.fill.solid()
    d.fill.fore_color.rgb = INK
    d.line.fill.background()
    d.shadow.inherit = False


def page_no(slide, n):
    tb, tf = textbox(slide, SW - 1.1, SH - 0.55, 0.7, 0.35)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, str(n), font=SANS, size=12, color=GRAY)


def running_title(slide, lines=("Neural", "Network")):
    """Two-line serif section header on the left."""
    tb, tf = textbox(slide, 0.62, 0.88, 3.5, 1.7)
    for i, ln in enumerate(lines):
        p = para(tf, first=(i == 0), space_after=0, line=0.98)
        run(p, ln, font=SERIF, size=40, color=INK, bold=True)


def fit_image(slide, path, box_x, box_y, box_w, box_h):
    """Place image scaled to fit inside the box, centred."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w
        h = box_w / ar
    else:
        h = box_h
        w = box_h * ar
    x = box_x + (box_w - w) / 2
    y = box_y + (box_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def caption(slide, text_bold, text_rest, cx, cy, cw):
    tb, tf = textbox(slide, cx, cy, cw, 0.6)
    p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.0)
    run(p, text_bold, font=SANS, size=11.5, color=INK, bold=True)
    run(p, text_rest, font=SANS, size=11.5, color=GRAY)


def subtitle(slide, text, x=0.62, y=2.55, w=3.7):
    tb, tf = textbox(slide, x, y, w, 1.0)
    p = para(tf, first=True, line=1.02)
    run(p, text, font=SANS, size=17, color=INK, bold=True)
    return tb


def bullets(slide, items, x=0.62, y=3.45, w=3.7, h=3.6, size=13.5, gap=9):
    tb, tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap, line=1.06)
        run(p, "•  ", font=SANS, size=size, color=GOLD, bold=True)   # gold bullet dot
        if isinstance(it, tuple):
            head, rest = it
            run(p, head, font=SANS, size=size, color=INK, bold=True)
            run(p, rest, font=SANS, size=size, color=INK)
        else:
            run(p, it, font=SANS, size=size, color=INK)
    return tb


def chip(slide, x, y, w, label, value, value_color=GOLD):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = CHIPBG
    box.line.color.rgb = CHIPLINE
    box.line.width = Pt(1)
    box.shadow.inherit = False
    box.adjustments[0] = 0.12
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    p = para(tf, first=True, align=PP_ALIGN.CENTER, space_after=2, line=1.0)
    run(p, value, font=SERIF, size=23, color=value_color, bold=True)
    p2 = para(tf, align=PP_ALIGN.CENTER, line=1.0)
    run(p2, label, font=SANS, size=10, color=GRAY)


# ===========================================================================
#  SLIDES
# ===========================================================================
page = START_PAGE


def content_slide(sub, items, fig=None, cap=None, bullet_size=13.5,
                  fig_box=(4.55, 1.25, 8.3, 5.05), bullet_y=3.45, bullet_h=3.6):
    """Standard left-header + bullets + right-figure slide."""
    global page
    s = add_slide()
    dash(s)
    running_title(s)
    subtitle(s, sub)
    bullets(s, items, size=bullet_size, y=bullet_y, h=bullet_h)
    if fig:
        bx, by, bw, bh = fig_box
        cap_h = 0.5 if cap else 0.0
        fit_image(s, FIG / fig, bx, by, bw, bh - cap_h)
        if cap:
            caption(s, cap[0], cap[1], bx, by + bh - cap_h + 0.05, bw)
    page_no(s, page)
    page += 1
    return s


# --- Slide 1 : section divider ---------------------------------------------
s = add_slide()
tb, tf = textbox(s, 1.2, 2.7, SW - 2.4, 2.1, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.05)
run(p, "Attributing Texts to an Author", font=SERIF, size=46, color=INK, bold=True)
p2 = para(tf, align=PP_ALIGN.CENTER, line=1.05, space_before=4)
run(p2, "using a Neural Network", font=SERIF, size=46, color=INK, bold=True)
page_no(s, page); page += 1

# --- Slide 2 : the task / setup (with mini architecture diagram) -----------
s = add_slide()
dash(s)
running_title(s)
subtitle(s, "The Task: Open-Set Attribution")
bullets(s, [
    ("5 known authors + a 6th “none” class", " — open-set recognition, not forced choice"),
    ("873 passages", ", each described by 107 stylometric features"),
    ("Neural network", " — ReLU hidden units, softmax over 6 classes"),
    ("Adam optimizer", ", cross-entropy loss, early stopping"),
    ("60 / 20 / 20 split", " — train / validation / test, stratified by class"),
], size=14, y=3.45, h=3.5)

# mini architecture diagram on the right: [107 features] -> [64 ReLU] -> [softmax 6]
boxes = [("Input", "107\nfeatures", RGBColor(0xE9, 0xF0, 0xF7), RGBColor(0x2E, 0x5A, 0x88)),
         ("Hidden (ReLU)", "64\nunits", RGBColor(0xF3, 0xEC, 0xF7), RGBColor(0x7A, 0x4E, 0x9C)),
         ("Softmax", "6\nclasses", RGBColor(0xEA, 0xF5, 0xEE), GREEN)]
bx0, by0, bw0, bh0 = 5.4, 3.05, 1.85, 1.55
gap0 = 0.85
title_b, ttf = textbox(s, 4.95, 2.0, 7.6, 0.7)
pp = para(ttf, first=True, align=PP_ALIGN.CENTER)
run(pp, "Network architecture", font=SANS, size=14, color=INK, bold=True)
for i, (lab, val, fillc, edgec) in enumerate(boxes):
    x = bx0 + i * (bw0 + gap0)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(by0),
                           Inches(bw0), Inches(bh0))
    b.fill.solid(); b.fill.fore_color.rgb = fillc
    b.line.color.rgb = edgec; b.line.width = Pt(1.5)
    b.shadow.inherit = False; b.adjustments[0] = 0.1
    tf2 = b.text_frame; tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    pv = para(tf2, first=True, align=PP_ALIGN.CENTER, line=1.0, space_after=2)
    run(pv, val, font=SERIF, size=20, color=edgec, bold=True)
    pl = para(tf2, align=PP_ALIGN.CENTER, line=1.0)
    run(pl, lab, font=SANS, size=10.5, color=GRAY)
    if i < len(boxes) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(x + bw0 + 0.12), Inches(by0 + bh0 / 2 - 0.13),
                                Inches(gap0 - 0.24), Inches(0.26))
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xB8, 0xB5, 0xAD)
        ar.line.fill.background(); ar.shadow.inherit = False
caption(s, "Grid search ", "explores 4 depths (1, 3, 10, 50 hidden layers) and 3 feature subsets.",
        4.95, 5.0, 7.6)
page_no(s, page); page += 1

# --- Slide 3 : feature selection (RF permutation importance) ---------------
content_slide(
    "Choosing the Features",
    [("Random Forest permutation importance", " ranks all 107 features"),
     ("Shuffle one feature", ", measure the drop in held-out accuracy"),
     ("Multivariate", " — credits the feature the forest truly relies on, not redundant copies"),
     ("Top drivers", ": punctuation rates, Honoré's R, function words"),
     ("Top-30 / top-50 subsets", " feed the grid search")],
    fig="rf_permutation_importance.png",
    cap=("Figure 1. ", "Most informative stylometric features (permutation importance)."),
    bullet_size=13.5,
    fig_box=(4.7, 1.55, 8.1, 4.7),
    bullet_y=3.2, bullet_h=3.8,
)

# --- Slide 4 : grid search / model selection -------------------------------
content_slide(
    "Grid Search & Model Selection",
    [("12 configurations", " — 3 feature subsets (top 30 / 50 / all 107) × 4 depths"),
     ("Ranked on validation", ", not the training set"),
     ("Winner retrained", " on train + validation, then tested once"),
     ("Test set touched once", " — an honest generalization estimate"),
     ("Best:", " all 107 features, a single 64-unit hidden layer (validation 0.943)")],
    fig="validation_selection.png",
    cap=("Figure 2. ", "Validation accuracy for every feature-subset × depth combination."),
    bullet_size=13.5,
    fig_box=(4.7, 1.7, 8.1, 4.4),
)

# --- Slide 5 : results (metric chips + confusion) --------------------------
s = add_slide()
dash(s)
running_title(s)
subtitle(s, "Results — Open Set (Six Classes)")
chip(s, 0.62, 3.05, 1.72, "Test accuracy", "94.9%")
chip(s, 2.46, 3.05, 1.72, "Weighted F1", "0.949")
chip(s, 0.62, 4.15, 1.72, "Classes", "6")
chip(s, 2.46, 4.15, 1.72, "None-class F1", "1.00", value_color=GREEN)
tb, tf = textbox(s, 0.62, 5.45, 3.6, 1.7)
p = para(tf, first=True, line=1.12)
run(p, "The ", font=SANS, size=12.5, color=INK)
run(p, "none class", font=SANS, size=12.5, color=GREEN, bold=True)
run(p, " is recovered perfectly — no unknown passage leaks into a known author. ",
    font=SANS, size=12.5, color=INK)
run(p, "The few errors", font=SANS, size=12.5, color=INK, bold=True)
run(p, " are minor confusions among the five (Raemon recall 0.88, Scott Alexander 0.90).",
    font=SANS, size=12.5, color=INK)
fit_image(s, FIG / "confusion.png", 5.0, 1.3, 7.6, 4.9)
caption(s, "Figure 3. ", "Six-class confusion matrix, including the none class.", 5.0, 6.25, 7.6)
page_no(s, page); page += 1

# --- Slide 6 : depth effect ------------------------------------------------
content_slide(
    "How Deep Should the Network Be?",
    [("Accuracy is flat", " from depth 1 to depth 10"),
     ("Depth 50 collapses", " to ≈ 0.17 — the 1-in-6 random baseline"),
     ("Not over-fitting", " — the deep net never starts learning"),
     ("Shallow wins", ": a single hidden layer is the right choice")],
    fig="depth_effect.png",
    cap=("Figure 4. ", "Validation accuracy versus network depth."),
    fig_box=(4.7, 1.9, 8.1, 3.9),
)

# --- Slide 7 : gradient diagnostic (table + figure) ------------------------
s = add_slide()
dash(s)
running_title(s)
subtitle(s, "Why Depth 50 Fails: Vanishing Gradients")
bullets(s, [
    ("Measured at initialization", " — Glorot-uniform, sklearn's ReLU scheme"),
    ("First-layer gradient", " collapses as the network deepens"),
    ("At depth 50", " it is ~8 orders of magnitude weaker"),
], size=12.5, y=3.5, h=1.6)

tbl_rows = [
    ("Depth", "1st-layer grad", "Ratio out/1st"),
    ("1", "2.0×10⁻²", "3.5"),
    ("3", "7.4×10⁻³", "4.2"),
    ("10", "4.2×10⁻⁴", "13.7"),
    ("50", "2.5×10⁻¹⁰", "2.2×10⁷"),
]
gx, gy = 0.62, 5.3
rowh = 0.33
colw = [0.9, 1.55, 1.1]
for ri, rowv in enumerate(tbl_rows):
    cx = gx
    for ci, cell in enumerate(rowv):
        cellbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(gy + ri * rowh),
                                     Inches(colw[ci]), Inches(rowh))
        cellbox.fill.solid()
        cellbox.fill.fore_color.rgb = INK if ri == 0 else CHIPBG
        cellbox.line.color.rgb = CHIPLINE
        cellbox.line.width = Pt(0.75)
        cellbox.shadow.inherit = False
        ctf = cellbox.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Inches(0.04); ctf.margin_right = Inches(0.02)
        ctf.margin_top = ctf.margin_bottom = 0
        cp = para(ctf, first=True, align=PP_ALIGN.CENTER, line=1.0)
        is_head = ri == 0
        is_bad = (ri == 4)
        run(cp, cell, font=SANS, size=10.5,
            color=(RGBColor(0xFF, 0xFF, 0xFF) if is_head else
                   (GOLD if (is_bad and ci > 0) else INK)),
            bold=is_head or (is_bad and ci == 0))
        cx += colw[ci]

fit_image(s, FIG / "gradient_health.png", 4.7, 1.7, 8.1, 4.2)
caption(s, "Figure 5. ", "Gradient RMS per layer (left) and at the input layer vs depth (right).",
        4.7, 5.95, 8.1)
page_no(s, page); page += 1

# --- Slide 8 : takeaways ---------------------------------------------------
s = add_slide()
dash(s)
running_title(s)
subtitle(s, "Takeaways", y=2.45)
bullets(s, [
    ("A shallow NN", " attributes passages across six classes at ~95% accuracy"),
    ("107 stylometric features", " carry enough signal — no deep network needed"),
    ("Random Forest permutation importance", " picks the most informative features"),
    ("Depth hurts", ": vanishing gradients stall the 50-layer net at initialization"),
    ("Open-set works", ": a dedicated “none” class is recovered perfectly")],
    size=15, y=3.2, w=11.8, h=3.8, gap=12)
page_no(s, page); page += 1

prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
